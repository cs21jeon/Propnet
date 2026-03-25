"""Proptalk APP 결제경로 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

os.chdir(os.path.dirname(os.path.abspath(__file__)))
os.chdir('..')  # Propnet root

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_step_number(slide, number, left, top):
    shape = slide.shapes.add_shape(9, left, top, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_arrow(slide, left, top):
    shape = slide.shapes.add_shape(13, left, top, Inches(0.6), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


def add_phone_image(slide, img_path, left, top, height):
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h
    actual_width = int(height * aspect)
    slide.shapes.add_picture(img_path, left, top, actual_width, height)
    return actual_width


def add_slide_content(slide, title_text, imgs_data, start_step, img_height=Inches(4.5)):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Section title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK

    # Calculate start x based on number of images
    n = len(imgs_data)
    if n == 3:
        x = Inches(0.7)
    else:
        x = Inches(1.5)

    for i, (path, title, desc) in enumerate(imgs_data):
        step_num = start_step + i
        add_step_number(slide, step_num, x + Inches(0.5), Inches(1.0))

        txBox = slide.shapes.add_textbox(x + Inches(1.1), Inches(0.95), Inches(3), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = DARK

        w = add_phone_image(slide, path, x, Inches(1.6), img_height)

        txBox = slide.shapes.add_textbox(x, Inches(6.3), Emu(w), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY

        x_end = x + Emu(w)

        if i < n - 1:
            if n == 3:
                add_arrow(slide, x_end + Inches(0.15), Inches(3.8))
                x = x_end + Inches(0.9)
            else:
                add_arrow(slide, x_end + Inches(0.3), Inches(3.8))
                x = x_end + Inches(1.2)


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = WHITE

txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
tf = txBox.text_frame

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Proptalk APP 결제경로"
run.font.size = Pt(40)
run.font.bold = True
run.font.color.rgb = DARK

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)
run2 = p2.add_run()
run2.text = "통화 녹음 AI 텍스트 변환 \u00b7 요약 서비스"
run2.font.size = Pt(20)
run2.font.color.rgb = GRAY

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(40)
run3 = p3.add_run()
run3.text = "\uc0c1\ud638: \ud504\ub86d\ub137 | \ub300\ud45c: \uc804\ucc3d\uc131 | \uc0ac\uc5c5\uc790\ub4f1\ub85d\ubc88\ud638: 308-33-01901"
run3.font.size = Pt(14)
run3.font.color.rgb = GRAY

p4 = tf.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "\uc0c1\uc810\uc544\uc774\ub514(MID): gpropth85u"
run4.font.size = Pt(14)
run4.font.color.rgb = GRAY

p5 = tf.add_paragraph()
p5.alignment = PP_ALIGN.CENTER
run5 = p5.add_run()
run5.text = "https://goldenrabbit.biz/proptalk/"
run5.font.size = Pt(14)
run5.font.color.rgb = BLUE

# ── Slide 2: Steps 1~3 ──
CAP = "proptalk/images/Capture/"
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_content(slide, "STEP 1~3: \uc571 \ub0b4 \uacb0\uc81c \uc9c4\uc785 \uacbd\ub85c", [
    (CAP + "Proptalk_screenshot_\ub85c\uadf8\uc778\ud654\uba74.jpg",
     "\u2460 \uc571 \ub85c\uadf8\uc778", "Google \uacc4\uc815\uc73c\ub85c\n\ub85c\uadf8\uc778"),
    (CAP + "Proptalk_screenshot_\uba54\uc778\ud654\uba74.jpg",
     "\u2461 \uba54\uc778\ud654\uba74", "\ud504\ub85c\ud544 \uba54\ub274 \uc9c4\uc785"),
    (CAP + "Proptalk_screenshot_\ud504\ub85c\ud544_\uc0c1\ub2e8.jpg",
     "\u2462 \ud504\ub85c\ud544 > \uad6c\ub3c5/\uacb0\uc81c", "\ucda9\uc804/\uc694\uae08\uc81c \ud130\uce58"),
], start_step=1)

# ── Slide 3: Steps 4~5 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_content(slide, "STEP 4~5: \uc6f9 \uacb0\uc81c \ud398\uc774\uc9c0 \uc9c4\uc785", [
    (CAP + "Proptalk_screenshot_\uad6c\ub3c5\uacb0\uc81c_\uc0ac\uc6a9\uc774\ub825.jpg",
     "\u2463 \uad6c\ub3c5/\uacb0\uc81c \uc0c1\uc138", "\uc6f9 \uacb0\uc81c URL \uc548\ub0b4\n\ube0c\ub77c\uc6b0\uc800\ub85c \uc774\ub3d9"),
    ("proptalk/images/\uacb0\uc81c\ud398\uc774\uc9c0.jpg",
     "\u2464 \uc694\uae08\uc81c \uc120\ud0dd (\uc6f9)", "\uc2dc\uac04\ud329 \ub610\ub294 \uad6c\ub3c5\n\uc694\uae08\uc81c \uc120\ud0dd \ud6c4 \uad6c\ub9e4"),
], start_step=4)

# ── Slide 4: Steps 6~7 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_content(slide, "STEP 6~7: \ud1a0\uc2a4\ud398\uc774\uba3c\uce20 \uacb0\uc81c \ubc0f \uc644\ub8cc", [
    ("proptalk/images/\uacb0\uc81c\ud654\uba74.jpg",
     "\u2465 \ud1a0\uc2a4 \uacb0\uc81c\ucc3d", "\uce74\ub4dc \uc120\ud0dd \ubc0f\n\uacb0\uc81c \uc815\ubcf4 \uc785\ub825"),
    ("proptalk/images/\uacb0\uc81c\ud654\uba7402.jpg",
     "\u2466 \uacb0\uc81c \uc644\ub8cc", "\uacb0\uc81c \uc644\ub8cc \ud655\uc778\n\uc794\uc5ec \uc2dc\uac04 \ucda9\uc804"),
], start_step=6)

output = "assets/\ud504\ub86d\ub137_Proptalk_APP_\uacb0\uc81c\uacbd\ub85c.pptx"
prs.save(output)
print(f"PPT saved: {output}")
